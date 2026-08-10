import os
import re
import io
import base64
import logging
from enum import Enum
from typing import Optional, List
from PIL import Image
import fitz  # PyMuPDF

try:
    from .image_preprocessing import preprocessing_service
except ImportError:
    try:
        from app.services.image_preprocessing import preprocessing_service
    except ImportError:
        preprocessing_service = type("MockPreprocessing", (), {
            "deskew": lambda self, img: img,
            "enhance_contrast": lambda self, img: img,
        })()

logger = logging.getLogger(__name__)

def _get_tesseract_cmd() -> str:
    """Locate Tesseract binary on Windows, macOS, or Linux, or return 'tesseract' default."""
    env_path = os.getenv("TESSERACT_PATH") or os.getenv("TESSERACT_CMD")
    if env_path and os.path.exists(env_path):
        return env_path

    # Common Windows installation locations
    win_paths = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        os.path.expandvars(r"%LOCALAPPDATA%\Programs\Tesseract-OCR\tesseract.exe"),
    ]
    for p in win_paths:
        if os.path.exists(p):
            return p

    # Common macOS / Linux locations
    unix_paths = [
        "/opt/homebrew/bin/tesseract",
        "/usr/local/bin/tesseract",
        "/usr/bin/tesseract",
    ]
    for p in unix_paths:
        if os.path.exists(p):
            return p

    return "tesseract"


try:
    import pytesseract
    TESSERACT_AVAILABLE = True
    tess_cmd = _get_tesseract_cmd()
    if os.path.exists(tess_cmd):
        pytesseract.pytesseract.tesseract_cmd = tess_cmd
except ImportError:
    TESSERACT_AVAILABLE = False


class FileType(str, Enum):
    DIGITAL_PDF = "digital_pdf"
    SCANNED_PDF = "scanned_pdf"
    IMAGE = "image"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    TXT = "txt"
    UNKNOWN = "unknown"


def _is_image_ext(ext: str) -> bool:
    return ext.lower() in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp", ".gif"}


def detect_file_type(file_path: str) -> FileType:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        try:
            doc = fitz.open(file_path)
            total_pages = len(doc)
            text_pages = 0
            for page in doc:
                text = page.get_text().strip()
                if len(text) > 30:
                    text_pages += 1
            doc.close()
            if text_pages > 0 and (text_pages / max(total_pages, 1)) > 0.3:
                return FileType.DIGITAL_PDF
            return FileType.SCANNED_PDF
        except Exception as e:
            logger.warning(f"Failed to inspect PDF {file_path}: {e}")
            return FileType.SCANNED_PDF

    if _is_image_ext(ext):
        return FileType.IMAGE
    if ext == ".docx":
        return FileType.DOCX
    if ext in {".xlsx", ".xls"}:
        return FileType.XLSX
    if ext in {".pptx", ".ppt"}:
        return FileType.PPTX
    if ext in {".txt", ".csv", ".json", ".xml", ".html", ".md"}:
        return FileType.TXT

    logger.warning(f"Unknown file type for {file_path}, treating as scanned")
    return FileType.SCANNED_PDF


def _extract_digital_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    pages = []
    for page_num, page in enumerate(doc, 1):
        blocks = page.get_text("dict").get("blocks", [])
        page_lines = []
        for block in blocks:
            btype = block.get("type", 0)
            if btype == 0:
                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    if not spans:
                        continue
                    text = "".join(s.get("text", "") for s in spans)
                    font_size = max(s.get("size", 12) for s in spans)
                    flags = spans[0].get("flags", 0)
                    is_bold = bool(flags & 16) or "bold" in spans[0].get("font", "").lower()

                    if is_bold and font_size > 11:
                        text = f"## {text.strip()}"
                    page_lines.append(text.strip())
            elif btype == 1:
                page_lines.append("[IMAGE]")
        page_text = "\n".join(line for line in page_lines if line)
        pages.append(f"<!-- Page {page_num} -->\n{page_text}")
    doc.close()
    result = "\n\n".join(pages)
    result = _normalize_markdown(result)
    return result


def _extract_docx(file_path: str) -> str:
    import docx
    doc = docx.Document(file_path)
    parts = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        style_name = (p.style.name or "").lower()
        if "heading 1" in style_name:
            parts.append(f"# {t}")
        elif "heading 2" in style_name:
            parts.append(f"## {t}")
        elif "heading 3" in style_name:
            parts.append(f"### {t}")
        else:
            parts.append(t)

    for table in doc.tables:
        rows_data = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows_data.append(cells)
        if rows_data:
            ncols = max(len(r) for r in rows_data)
            sep = "| " + " | ".join(["---"] * ncols) + " |"
            parts.append("| " + " | ".join(rows_data[0]) + " |")
            parts.append(sep)
            for cells in rows_data[1:]:
                parts.append("| " + " | ".join(cells) + " |")
            parts.append("")

    result = "\n\n".join(parts)
    return _normalize_markdown(result)


def _extract_txt(file_path: str) -> str:
    encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    with open(file_path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


def _extract_xlsx(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}")
        rows_data = []
        for row in ws.iter_rows(values_only=True):
            if not any(row):
                continue
            cells = [str(c).strip() if c is not None else "" for c in row]
            rows_data.append(cells)
        if rows_data:
            ncols = max(len(r) for r in rows_data)
            sep = "| " + " | ".join(["---"] * ncols) + " |"
            parts.append("| " + " | ".join(rows_data[0]) + " |")
            parts.append(sep)
            for cells in rows_data[1:]:
                parts.append("| " + " | ".join(cells) + " |")
        parts.append("")
    wb.close()
    result = "\n".join(parts)
    return _normalize_markdown(result)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- Slide {slide_idx} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                table = shape.table
                md_rows = []
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    md_rows.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                slide_texts.append("\n".join(md_rows))
        parts.append("\n".join(slide_texts))
    result = "\n\n".join(parts)
    return _normalize_markdown(result)


def _preprocess_image(img: Image.Image) -> Image.Image:
    try:
        if img.mode != "RGB":
            img = img.convert("RGB")
        img = preprocessing_service.deskew(img)
        img = preprocessing_service.enhance_contrast(img)
        return img
    except Exception as e:
        logger.warning(f"Image preprocessing failed: {e}")
        return img


def _page_to_image(page) -> tuple[str, Image.Image]:
    pix = page.get_pixmap(dpi=300, alpha=False)
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    target_w = 1600
    w_percent = target_w / float(img.width)
    h_size = int(float(img.height) * float(w_percent))
    img_resized = img.resize((target_w, h_size), Image.LANCZOS)
    buf = io.BytesIO()
    img_resized.save(buf, format="JPEG", quality=90, optimize=True)
    return base64.b64encode(buf.getvalue()).decode("utf-8"), img


def _load_image_b64(file_path: str) -> tuple[str, Image.Image]:
    img = Image.open(file_path)
    img = _preprocess_image(img)
    target_w = 1600
    w_percent = target_w / float(img.width)
    h_size = int(float(img.height) * float(w_percent))
    img_resized = img.resize((target_w, h_size), Image.LANCZOS)
    buf = io.BytesIO()
    img_resized.save(buf, format="JPEG", quality=90, optimize=True)
    return base64.b64encode(buf.getvalue()).decode("utf-8"), img


VISION_PROMPT = """You are a precision OCR engine and document data extractor.

Instructions:
1. Extract ALL visible text, numbers, tables, headings, customer details, invoice metadata, and labels from this image with 100% fidelity.
2. Format ALL tabular data strictly as Markdown tables (| Col 1 | Col 2 | ... |). Extract EVERY SINGLE row (e.g., items 1, 2, ... 11, 12, 13) including product descriptions, SKUs, quantities, rates, and amounts.
3. If a page contains a continuation table (e.g., items 11, 12, 13), format those continuation rows strictly as a Markdown table as well.
4. Output ONLY clean Markdown text starting directly with the content.
5. ABSOLUTELY NO internal thinking, reasoning out loud, self-corrections, or phrases like "Wait, looking at..." or "Let me re-examine...". Output ONLY the visible document content."""


def _strip_think_tags(text: str) -> str:
    if not text:
        return ""
    cleaned = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
    if "</think>" in cleaned:
        cleaned = cleaned.split("</think>")[-1]

    # Strip Vision LLM internal monologue and repetitive chain-of-thought reasoning loops
    reasoning_keywords = [
        "wait, looking", "let me re-examine", "let's look at",
        "crop ", "transcribe what is visible", "looking closely at",
        "looking at the very bottom", "i can see `sku:", "let's assume the sku"
    ]
    filtered_lines = []
    for line in cleaned.splitlines():
        lower = line.lower().strip()
        if any(kw in lower for kw in reasoning_keywords):
            continue
        filtered_lines.append(line)

    return "\n".join(filtered_lines).strip()


def _extract_tesseract_page_text(img: Image.Image) -> str:
    """Run local Tesseract; prefer spatial layout reconstruction for invoice tables."""
    spatial = _tesseract_spatial_markdown(img)
    if spatial and spatial.count("|") >= 6:
        return spatial
    try:
        import subprocess, tempfile
        tess_bin = _get_tesseract_cmd()
        tmp_img = tempfile.NamedTemporaryFile(suffix=".png", delete=False).name
        img.save(tmp_img)
        proc = subprocess.run(
            [tess_bin, tmp_img, "stdout", "--psm", "6", "-l", "eng"],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True
        )
        if os.path.exists(tmp_img):
            os.remove(tmp_img)
        plain = (proc.stdout or "").strip()
        if spatial and len(spatial) > len(plain) * 0.5:
            return f"{plain}\n\n{spatial}" if plain else spatial
        return plain
    except Exception:
        return spatial or ""


def _parse_money_token(tok: str) -> float | None:
    t = tok.replace("Rs.", "").replace("Rs", "").replace("PKR", "").replace(",", "").strip()
    if not re.fullmatch(r"\d+(\.\d+)?", t):
        return None
    try:
        return float(t)
    except ValueError:
        return None


def _tesseract_spatial_markdown(img: Image.Image) -> str:
    """
    Rebuild invoice-like tables using word bounding boxes so Qty / Rate / Amount
    stay aligned (plain image_to_string loses Digilog-style right-side columns).
    """
    if not TESSERACT_AVAILABLE:
        return ""
    try:
        import pytesseract

        tess_bin = _get_tesseract_cmd()
        if os.path.exists(tess_bin):
            pytesseract.pytesseract.tesseract_cmd = tess_bin

        data = pytesseract.image_to_data(
            img, output_type=pytesseract.Output.DICT, config="--psm 6"
        )
        width, _height = img.size
        words: list[dict] = []
        n = len(data.get("text") or [])
        for i in range(n):
            t = (data["text"][i] or "").strip()
            if not t:
                continue
            conf_raw = str(data["conf"][i])
            try:
                conf = int(float(conf_raw))
            except ValueError:
                conf = -1
            if conf < 0:
                continue
            words.append(
                {
                    "t": t,
                    "l": int(data["left"][i]),
                    "top": int(data["top"][i]),
                    "h": max(int(data["height"][i]), 1),
                }
            )
        if len(words) < 8:
            return ""

        words.sort(key=lambda w: (w["top"], w["l"]))
        rows: list[dict] = []
        for w in words:
            placed = False
            for row in rows:
                if abs(row["top"] - w["top"]) <= max(12, w["h"] * 0.65):
                    row["words"].append(w)
                    nwords = len(row["words"])
                    row["top"] = (row["top"] * (nwords - 1) + w["top"]) / nwords
                    placed = True
                    break
            if not placed:
                rows.append({"top": float(w["top"]), "words": [w]})

        desc_max_x = width * 0.50
        qty_min_x = width * 0.48
        money_min_x = width * 0.58

        def pick_qty_rate_amount(row_words: list[dict]) -> tuple[float | None, float | None, float | None]:
            qty_cands: list[float] = []
            money_cands: list[float] = []
            for w in sorted(row_words, key=lambda x: x["l"]):
                val = _parse_money_token(w["t"])
                if val is None:
                    continue
                # Skip obvious invoice/order ids
                if val >= 100000:
                    continue
                if w["l"] >= money_min_x:
                    money_cands.append(val)
                elif qty_min_x <= w["l"] < money_min_x and val == int(val) and 1 <= val <= 5000:
                    qty_cands.append(val)

            # Digilog pattern on the right: Qty, Rate, Amount all in money zone
            # e.g. 25, 130, 3250
            nums = []
            for w in sorted(row_words, key=lambda x: x["l"]):
                if w["l"] < qty_min_x:
                    continue
                val = _parse_money_token(w["t"])
                if val is None or val >= 100000:
                    continue
                nums.append(val)

            def score(q, r, a) -> float:
                if q is None or r is None or a is None or q <= 0 or r <= 0:
                    return 1e18
                return abs(q * r - a)

            best = (None, None, None)
            best_s = 1e18

            # Try consecutive triples from the right-side number stream
            for i in range(len(nums)):
                for j in range(i + 1, len(nums)):
                    for k in range(j + 1, len(nums)):
                        for q, r, a in (
                            (nums[i], nums[j], nums[k]),
                            (nums[j], nums[i], nums[k]),
                        ):
                            s = score(q, r, a)
                            # Prefer small integer qty
                            if q != int(q) or q > 5000:
                                s += 1000
                            if s < best_s:
                                best_s = s
                                best = (float(q), float(r), float(a))

            if best[0] is not None and best_s <= max(1.0, 0.05 * (best[2] or 1)):
                return best

            if len(nums) >= 2:
                # last two are rate + amount; qty from qty band or amount/rate
                rate, amount = nums[-2], nums[-1]
                qty = qty_cands[0] if qty_cands else None
                if qty is None and rate > 0:
                    q2 = amount / rate
                    if abs(q2 - round(q2)) < 0.08 and round(q2) >= 1:
                        qty = float(round(q2))
                    else:
                        qty = 1.0
                if score(qty, rate, amount) <= max(1.0, 0.05 * amount):
                    return qty, rate, amount
                return qty or 1.0, rate, amount

            if len(nums) == 1:
                return 1.0, nums[0], nums[0]
            return None, None, None

        items: list[dict] = []
        header_lines: list[str] = []

        def flush_pending_desc(text: str):
            if not items or not text:
                return
            prev = items[-1]
            if not prev.get("sku") and re.search(r"\bSKU\s*:", text, re.I):
                prev["sku"] = re.sub(r"^.*?\bSKU\s*:\s*", "", text, flags=re.I).strip()
            elif not re.search(r"\bSKU\s*:", text, re.I):
                # continuation of description
                if not re.search(r"@|madina|faisalabad|invoice\s*date|order\s*number", text, re.I):
                    prev["description"] = f"{prev['description']} {text}".strip()

        junk_desc = re.compile(
            r"invoice\s*date|order\s*number|order\s*date|payment\s*method|shipping\s*method|"
            r"madina\s*town|faisalabad|@gmail|subtotal|grand\s*total|nigilog|electronics\b",
            re.I,
        )

        for row in rows:
            ws = sorted(row["words"], key=lambda x: x["l"])
            left = [w for w in ws if w["l"] < desc_max_x]
            left_text = " ".join(w["t"] for w in left).strip()
            left_text = re.sub(r"^[\d\.\)\|]+\s*", "", left_text)
            left_text = re.sub(r"^[^A-Za-z0-9]+", "", left_text).strip()

            qty, rate, amount = pick_qty_rate_amount(ws)
            full_line = " ".join(w["t"] for w in ws).strip()

            if amount is not None and left_text and len(left_text) > 3 and not junk_desc.search(left_text):
                # Reject header rows that swallowed order id as amount
                if amount >= 100000 or (
                    qty is None and rate is None and junk_desc.search(full_line or "")
                ):
                    if full_line:
                        header_lines.append(full_line)
                    continue
                items.append(
                    {
                        "description": left_text,
                        "sku": None,
                        "quantity": qty,
                        "unit_price": rate,
                        "total_price": amount,
                    }
                )
            elif left_text and re.search(r"\bSKU\s*:", left_text, re.I):
                flush_pending_desc(left_text)
            elif left_text and items and amount is None:
                flush_pending_desc(left_text)
            elif full_line and not items:
                header_lines.append(full_line)
            elif left_text and not items:
                header_lines.append(left_text)

        if len(items) < 2:
            return ""

        lines = header_lines[:20]
        lines.append("")
        lines.append("### Line items (spatial OCR)")
        lines.append("| # | Description | SKU | Qty | Unit price | Amount |")
        lines.append("|---|---|---|---:|---:|---:|")
        for i, it in enumerate(items, 1):
            desc = (it["description"] or "").replace("|", "/")
            sku = (it.get("sku") or "").replace("|", "/")
            q = it.get("quantity")
            r = it.get("unit_price")
            a = it.get("total_price")
            q_s = f"{q:g}" if isinstance(q, (int, float)) else ""
            r_s = f"{r:g}" if isinstance(r, (int, float)) else ""
            a_s = f"{a:g}" if isinstance(a, (int, float)) else ""
            lines.append(f"| {i} | {desc} | {sku} | {q_s} | {r_s} | {a_s} |")

        total_sum = sum(
            float(it["total_price"])
            for it in items
            if isinstance(it.get("total_price"), (int, float))
        )
        if total_sum > 0:
            lines.append("")
            lines.append(f"Spatial OCR line sum: {total_sum:g}")

        return "\n".join(lines)
    except Exception as e:
        logger.warning(f"Spatial Tesseract failed: {e}")
        return ""


def _merge_tesseract_header(page_text: str, tess_text: str) -> str:
    """Generically prepend pre-table header lines from Tesseract if missing from Vision OCR output."""
    if not tess_text:
        return page_text

    lower_vision = page_text.lower()
    header_lines = []

    for line in tess_text.splitlines():
        l_str = line.strip()
        if not l_str:
            continue
        # Stop header extraction as soon as table rows or table headers start
        if l_str.startswith("|") or l_str.startswith("#") or re.match(r"^\d+[\.\)]\s+", l_str):
            break

        # Check if this non-table line is already in page_text
        words = [w for w in re.findall(r"\w+", l_str.lower()) if len(w) > 2]
        already_present = any(w in lower_vision for w in words[:3]) if words else False
        if not already_present:
            header_lines.append(l_str)

    if header_lines:
        header_block = "\n".join(header_lines)
        return f"{header_block}\n\n{page_text}"

    return page_text


def _format_tesseract_items_table(text: str) -> str:
    """Generically format un-tabled structured text into Markdown tables if Vision API is unavailable."""
    if not text or "|" in text:
        return text

    import re
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    if not lines:
        return text

    # Detect lines containing numeric item prefixes, currency symbols, quantities, or SKU/ID patterns
    item_pattern = re.compile(
        r"^\d+[\.\)]\s+|"  # Item number e.g. "1.", "1)"
        r"\b(?:rs\.?|pkr|\$|eur|usd|gbp|sku|code|part\s*#?|qty|pcs)\b|"  # Financial/Inventory identifiers
        r"\b\d+(?:\.\d{1,2})?\s*(?:rs\.?|pkr|\$|eur|usd)\b",
        re.IGNORECASE
    )

    table_rows = [l for l in lines if item_pattern.search(l)]

    if len(table_rows) >= 2:
        md_lines = [
            "| # | Product / Item Description | Details / Amounts |",
            "|---|---|---|",
        ]
        for idx, row in enumerate(table_rows, 1):
            md_lines.append(f"| {idx} | {row} | |")

        return text + "\n\n### Extracted Structured Data\n" + "\n".join(md_lines)

    return text


def _vision_ocr(image_b64s: list[str], pil_images: list[Image.Image] | None = None) -> str:
    """Hybrid per-page OCR: Vision first, instant Tesseract fallback per page."""
    from .groq_service import groq_service
    from .vision_provider import vision_provider

    page_texts = []
    for idx, b64 in enumerate(image_b64s, 1):
        content = [
            {"type": "text", "text": VISION_PROMPT},
            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}}
        ]
        page_text = ""
        tess_text = ""
        if pil_images and idx - 1 < len(pil_images):
            tess_text = _extract_tesseract_page_text(pil_images[idx - 1])

        # Attempt 1: Active Groq Vision
        if groq_service.available and groq_service.vision_available:
            try:
                result = groq_service.chat_vision(
                    [{"role": "user", "content": content}],
                    temperature=0.0,
                    max_tokens=4096,
                )
                if result and not result.startswith("[Groq") and not result.startswith("⚠️"):
                    page_text = result
            except Exception as e:
                print(f"[OCR] Groq vision page {idx} failed: {e}")

        # Attempt 2: Multi-provider Vision Fallback (Gemini, OpenAI, Anthropic)
        if not page_text.strip():
            try:
                res = vision_provider.analyze(b64)
                ocr_out = res.get("ocr_text", "") or res.get("markdown", "")
                if ocr_out and not ocr_out.startswith("["):
                    page_text = ocr_out
            except Exception as e:
                print(f"[OCR] Multi-provider vision failed: {e}")

        # Attempt 3: Local Tesseract OCR fallback for page
        if not page_text.strip():
            page_text = _format_tesseract_items_table(tess_text)
            if page_text:
                print(f"[OCR] Tesseract extracted {len(page_text)} chars for page {idx}")

        if page_text.strip():
            # Merge missing customer/vendor/invoice metadata from Tesseract if Vision skipped them
            if tess_text:
                page_text = _merge_tesseract_header(page_text, tess_text)

            page_text = _strip_think_tags(page_text)
            page_texts.append(f"<!-- Page {idx} -->\n{page_text.strip()}")

    return "\n\n".join(page_texts)


def _tesseract_ocr(images: list[Image.Image]) -> str:
    if not TESSERACT_AVAILABLE:
        logger.warning("Tesseract not installed, skipping fallback OCR")
        return ""
    try:
        import pytesseract

        tess_bin = _get_tesseract_cmd()
        if os.path.exists(tess_bin):
            pytesseract.pytesseract.tesseract_cmd = tess_bin
        texts = []
        for img in images:
            spatial = _tesseract_spatial_markdown(img)
            plain = pytesseract.image_to_string(img, config="--psm 6") or ""
            plain = plain.strip()
            if spatial and spatial.count("|") >= 6:
                # Prefer structured table; keep plain header if useful
                header = []
                for line in plain.splitlines():
                    s = line.strip()
                    if not s:
                        continue
                    if re.match(r"^\d+[\.\)]\s+", s) or s.startswith("|"):
                        break
                    header.append(s)
                block = "\n".join(header[:25])
                texts.append(f"{block}\n\n{spatial}".strip() if block else spatial)
            elif plain:
                texts.append(_format_tesseract_items_table(plain))
            elif spatial:
                texts.append(spatial)
        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"Tesseract OCR failed: {e}")
        return ""


def process_scanned_pdf(file_path: str) -> str:
    # Cap Vision OCR pages — sequential Groq Vision per page is the main latency driver.
    max_vision_pages = int(os.getenv("OCR_MAX_VISION_PAGES", "8"))
    images = []
    b64s = []
    try:
        doc = fitz.open(file_path)
        for i, page in enumerate(doc):
            if i >= max_vision_pages:
                break
            # Higher DPI helps Digilog-style scanned invoices keep Qty/Rate readable
            pix = page.get_pixmap(dpi=250, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            target_w = 1800
            w_percent = target_w / float(img.width)
            h_size = int(float(img.height) * float(w_percent))
            img = img.resize((target_w, h_size), Image.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=92, optimize=True)
            b64s.append(base64.b64encode(buf.getvalue()).decode("utf-8"))
            images.append(img)
        total_pages = doc.page_count
        doc.close()
        if total_pages > max_vision_pages:
            logger.info(
                f"[OCR] Limiting Vision OCR to first {max_vision_pages}/{total_pages} pages for speed"
            )
    except Exception as e:
        logger.error(f"Failed to load PDF pages for OCR: {e}")

    # Prefer spatial Tesseract for scanned invoices (column-aware Qty/Rate/Amount)
    if images and TESSERACT_AVAILABLE:
        try:
            tess_text = _tesseract_ocr(images)
            has_spatial_table = "### Line items (spatial OCR)" in (tess_text or "")
            if tess_text.strip() and (has_spatial_table or len(tess_text.strip()) >= 200):
                logger.info(
                    "Using Tesseract OCR%s — sufficient text extracted",
                    " (spatial tables)" if has_spatial_table else "",
                )
                return _normalize_markdown(tess_text)
        except Exception as e:
            logger.warning(f"Tesseract pre-pass failed: {e}")

    # 1. Hybrid Vision + Tesseract OCR (per-page fallback)
    if b64s:
        try:
            v_text = _vision_ocr(b64s, pil_images=images)
            if v_text.strip() and not v_text.startswith("["):
                # Append spatial tables if Vision missed Qty/Rate columns
                if images and TESSERACT_AVAILABLE and "### Line items (spatial OCR)" not in v_text:
                    try:
                        spatial_extra = _tesseract_ocr(images)
                        if "### Line items (spatial OCR)" in (spatial_extra or ""):
                            v_text = f"{v_text}\n\n{spatial_extra}"
                    except Exception:
                        pass
                return _normalize_markdown(v_text)
        except Exception as e:
            logger.warning(f"Vision OCR failed with exception: {e}, shifting to local OCR fallback")

    # 2. Fallback: Local Tesseract OCR
    if images:
        try:
            tess_text = _tesseract_ocr(images)
            if tess_text.strip():
                logger.info("Successfully extracted text via local Tesseract OCR fallback")
                return _normalize_markdown(tess_text)
        except Exception as e:
            logger.warning(f"Tesseract OCR fallback failed: {e}")

    # 3. Fallback: PyMuPDF Direct Text Extraction
    try:
        doc = fitz.open(file_path)
        parts = []
        for page_num, page in enumerate(doc, 1):
            t = page.get_text().strip()
            if t:
                parts.append(f"<!-- Page {page_num} -->\n{t}")
        doc.close()
        fallback = "\n\n".join(parts)
        if fallback.strip():
            logger.info("Successfully extracted text via PyMuPDF fallback")
            return _normalize_markdown(fallback)
    except Exception as e:
        logger.error(f"PyMuPDF fallback failed: {e}")

    return "[OCR failed: All Vision and Local OCR methods failed]"


def process_image(file_path: str) -> str:
    b64, img = None, None
    try:
        b64, img = _load_image_b64(file_path)
    except Exception as e:
        logger.error(f"Failed to load image for OCR: {e}")

    # 1. Attempt Vision OCR
    if b64:
        try:
            v_text = _vision_ocr([b64], pil_images=[img] if img else None)
            if v_text.strip() and not v_text.startswith("["):
                return _normalize_markdown(v_text)
        except Exception as e:
            logger.warning(f"Vision OCR failed for image: {e}, shifting to local OCR fallback")

    # 2. Fallback: Local Tesseract OCR
    if img:
        try:
            tess_text = _tesseract_ocr([img])
            if tess_text.strip():
                return _normalize_markdown(tess_text)
        except Exception as e:
            logger.warning(f"Tesseract OCR failed for image: {e}")

    return "[OCR failed: All Vision and Local OCR methods failed]"


def _normalize_markdown(text: str) -> str:
    text = re.sub(r'\n{3,}', '\n\n', text)
    lines = text.split('\n')
    result = []
    for line in lines:
        if line.startswith('#'):
            stripped = re.sub(r'^(#+)\s*', r'\1 ', line)
            result.append(stripped)
        else:
            result.append(line)
    return '\n'.join(result)


def process_document(file_path: str) -> dict:
    file_type = detect_file_type(file_path)
    logger.info(f"[OCR] Detected: {file_type.value} — {file_path}")

    text = ""
    page_count = 0
    source = "vision"

    if file_type == FileType.DIGITAL_PDF:
        text = _extract_digital_pdf(file_path)
        # Fast path: keep native PDF text when it is usable.
        # Do NOT fall into Vision OCR just because pages contain images —
        # that path is multi-minute and blocks extraction.
        if text.strip() and len(text.strip()) >= 120:
            source = "digital_pdf"
        elif not text.strip() or len(text.strip()) < 120:
            vision_text = process_scanned_pdf(file_path)
            if vision_text and vision_text != "[OCR failed]" and len(vision_text) > len(text):
                text = vision_text
                source = "vision"
            else:
                source = "digital_pdf"
        try:
            doc = fitz.open(file_path)
            page_count = doc.page_count
            doc.close()
        except Exception:
            page_count = 1

    elif file_type == FileType.DOCX:
        text = _extract_docx(file_path)
        page_count = max(1, len(text) // 2000)
        source = "docx"

    elif file_type == FileType.XLSX:
        text = _extract_xlsx(file_path)
        page_count = max(1, len(text) // 2000)
        source = "xlsx"

    elif file_type == FileType.PPTX:
        text = _extract_pptx(file_path)
        page_count = max(1, len(text) // 2000)
        source = "pptx"

    elif file_type == FileType.TXT:
        text = _extract_txt(file_path)
        page_count = max(1, len(text) // 2000)
        source = "txt"

    elif file_type == FileType.SCANNED_PDF:
        text = process_scanned_pdf(file_path)
        if text and text != "[OCR failed]":
            try:
                doc = fitz.open(file_path)
                page_count = doc.page_count
                doc.close()
            except Exception:
                page_count = 0
        source = "vision"

    elif file_type == FileType.IMAGE:
        text = process_image(file_path)
        page_count = 1
        source = "vision"

    logger.info(f"[OCR] Result: {len(text)} chars, {page_count} pages, source={source}")
    return {"text": text, "page_count": page_count, "source": source}


ocr_service = type("OCRService", (), {"process_document": staticmethod(process_document)})()
